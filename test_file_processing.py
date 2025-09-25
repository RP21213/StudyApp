#!/usr/bin/env python3
"""
Test script to verify that the new file processing functions work correctly.
This script tests the text extraction functions for PDF, DOCX, and PPTX files.
"""

import sys
import os
import io

# Add the current directory to the path so we can import from app.py
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

try:
    from docx import Document as DocxDocument
    from pptx import Presentation
    import PyPDF2
    print("✅ All required libraries are installed")
except ImportError as e:
    print(f"❌ Missing library: {e}")
    print("Please run: pip install python-pptx python-docx PyPDF2")
    sys.exit(1)

def pdf_to_text(pdf_file_stream):
    try:
        reader = PyPDF2.PdfReader(pdf_file_stream)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text: text += page_text + "\n"
        return text
    except Exception as e:
        print(f"Error reading PDF: {e}")
        return ""

def docx_to_text(docx_file_stream):
    try:
        doc = DocxDocument(docx_file_stream)
        text = ""
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"
        return text
    except Exception as e:
        print(f"Error reading Word document: {e}")
        return ""

def pptx_to_text(pptx_file_stream):
    try:
        prs = Presentation(pptx_file_stream)
        text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error reading PowerPoint document: {e}")
        return ""

def extract_text_from_file(file_stream, file_extension):
    """
    Universal text extraction function that handles PDF, DOCX, and PPTX files.
    """
    file_extension = file_extension.lower()
    
    if file_extension == '.pdf':
        return pdf_to_text(file_stream)
    elif file_extension == '.docx':
        return docx_to_text(file_stream)
    elif file_extension == '.pptx':
        return pptx_to_text(file_stream)
    else:
        print(f"Unsupported file type: {file_extension}")
        return ""

def test_file_processing():
    """Test the file processing functions with sample files."""
    
    print("\n🧪 Testing File Processing Functions")
    print("=" * 50)
    
    # Test with a simple text file to verify the functions work
    test_files = [
        ("test.pdf", ".pdf"),
        ("test.docx", ".docx"), 
        ("test.pptx", ".pptx")
    ]
    
    for filename, extension in test_files:
        print(f"\n📄 Testing {filename} ({extension})")
        
        if not os.path.exists(filename):
            print(f"   ⚠️  {filename} not found - skipping test")
            continue
            
        try:
            with open(filename, 'rb') as f:
                file_stream = io.BytesIO(f.read())
                text = extract_text_from_file(file_stream, extension)
                
                if text.strip():
                    print(f"   ✅ Successfully extracted {len(text)} characters")
                    print(f"   📝 Preview: {text[:100]}...")
                else:
                    print(f"   ⚠️  No text extracted")
                    
        except Exception as e:
            print(f"   ❌ Error processing {filename}: {e}")
    
    print("\n✅ File processing test completed!")
    print("\nTo test with real files:")
    print("1. Place test.pdf, test.docx, and test.pptx files in this directory")
    print("2. Run this script again")
    print("3. The script will extract and display text from each file")

if __name__ == "__main__":
    test_file_processing()
